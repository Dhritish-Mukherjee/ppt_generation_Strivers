"""
generate_pptx.py
----------------
Replaces question content in the template.pptx with new questions
from questions.json, preserving all original formatting.
Now includes title slide image replacement.

Usage:
    python generate_pptx.py
    python generate_pptx.py --questions my_questions.json --output output.pptx --image new_cover.jpg
"""

import json
import argparse
import re
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Formatting constants (matched from original template) ──────────────────
QUESTION_COLOR = RGBColor(0xFF, 0xFF, 0x00)   # Yellow
OPTION_COLOR   = RGBColor(0xFF, 0xFF, 0xFF)   # White
OPTION_BOLD    = True

# Font size tiers (in EMUs: 1pt = 12700 EMU)
SIZE_LARGE  = int(48 * 12700)   # ≤ 400 total chars
SIZE_MEDIUM = int(42 * 12700)   # 401–600 total chars
SIZE_SMALL  = int(35 * 12700)   # > 600 total chars


def pick_font_size(question):
    """
    Auto-select font size based on total content length.
    Counts all text: both question lines + all option lines.
    """
    all_text = (
        question['question_en'] + question['question_bn'] +
        "".join(o['en'] + o.get('bn', '') for o in question['options'])
    )
    total = len(all_text)
    if total <= 400:
        return SIZE_LARGE, int(SIZE_LARGE * 0.82)   # question, options
    elif total <= 600:
        return SIZE_MEDIUM, int(SIZE_MEDIUM * 0.82)
    else:
        return SIZE_SMALL, int(SIZE_SMALL * 0.82)


# ── Helpers ────────────────────────────────────────────────────────────────

def replace_title_image(slide, image_path):
    """
    Replaces the first picture found on the given slide with a new image.
    Manipulates XML to ensure the new image stays on the same layer (z-order).
    Aggressively checks for standard pictures and picture placeholders.
    """
    for shape in slide.shapes:
        is_standard_pic = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        is_placeholder_pic = (shape.is_placeholder and hasattr(shape, "image"))
        is_named_pic = ("Picture" in shape.name or "Image" in shape.name)

        if is_standard_pic or is_placeholder_pic or is_named_pic:
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                old_element = shape._element
                
                new_pic = slide.shapes.add_picture(image_path, left, top, width, height)
                new_element = new_pic._element
                
                old_element.addprevious(new_element)
                old_element.getparent().remove(old_element)
                
                print(f"  ✓ Title image successfully replaced with: {image_path}")
                return
            except Exception as e:
                continue
            
    print("  ⚠ No picture shape was found on the title slide to replace.")
    print("    (Tip: If your template image is set as the 'Slide Background', open the template,")
    print("     remove the background, paste the image normally, and send to back!)")

def find_question_shape(slide):
    """Return the shape that contains the question using a strict regex pattern."""
    pattern = re.compile(r'^\d+\.\s*')
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if pattern.match(text):
                return shape
    return None


def clear_extra_runs(para, keep=1):
    """Remove extra runs from a paragraph, keeping only the first `keep` runs."""
    p_elem = para._p
    runs = p_elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r')
    for run in runs[keep:]:
        p_elem.remove(run)


def safe_set_run(para, text, size, color, bold=None):
    """Safely updates a paragraph's text, injecting a run if missing to prevent silent failures."""
    if not para.runs:
        para.add_run()
    
    clear_extra_runs(para, keep=1)
    run = para.runs[0]
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def replace_question_content(shape, q_num, question):
    """
    Replace text in the question textbox dynamically mapping paragraphs
    to ignore boilerplate text and extra line breaks.
    """
    tf = shape.text_frame
    paras = tf.paragraphs
    q_size, opt_size = pick_font_size(question)
    
    para_map = {
        'Q_en': 0,
        'Q_bn': 1,
        'opts': {}
    }
    
    for idx, p in enumerate(paras):
        text = p.text.strip()
        if text.startswith('A.'): para_map['opts']['A'] = idx
        elif text.startswith('B.'): para_map['opts']['B'] = idx
        elif text.startswith('C.'): para_map['opts']['C'] = idx
        elif text.startswith('D.'): para_map['opts']['D'] = idx

    q_text_en = f"{q_num}. {question['question_en']}"
    if len(paras) > 0:
        safe_set_run(paras[0], q_text_en, q_size, QUESTION_COLOR)
        
    if len(paras) > 1:
        safe_set_run(paras[1], question['question_bn'], q_size, QUESTION_COLOR)

    prefixes = ['A', 'B', 'C', 'D']
    
    is_compact = ('A' in para_map['opts'] and 'B' in para_map['opts'] and 
                  para_map['opts']['B'] == para_map['opts']['A'] + 1)

    for i, opt_data in enumerate(question['options']):
        if i >= len(prefixes): 
            break
            
        letter = prefixes[i]
        
        clean_en = re.sub(r'^[A-D]\.\s*', '', opt_data['en']).strip()
        clean_bn = re.sub(r'^[A-D]\.\s*', '', opt_data.get('bn', '')).strip()
        
        if letter in para_map['opts']:
            target_idx = para_map['opts'][letter]
            target_para = paras[target_idx]
            
            if is_compact:
                combined = f"{letter}. {clean_en} {clean_bn}".strip()
                safe_set_run(target_para, combined, opt_size, OPTION_COLOR, OPTION_BOLD)
            else:
                en_text = f"{letter}. {clean_en}"
                safe_set_run(target_para, en_text, opt_size, OPTION_COLOR, OPTION_BOLD)
                
                bn_idx = target_idx + 1
                
                if bn_idx < len(paras) and bn_idx not in para_map['opts'].values():
                    safe_set_run(paras[bn_idx], clean_bn, opt_size, OPTION_COLOR, OPTION_BOLD)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Generate quiz PPTX from template")
    parser.add_argument('--template',  default='template.pptx',  help='Path to template.pptx')
    parser.add_argument('--questions', default='questions.json', help='Path to questions JSON file')
    parser.add_argument('--output',    default='output.pptx',    help='Output file name')
    parser.add_argument('--image',     default=None,             help='Path to a new image for the Title Slide (Slide 1)')
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print(f"❌ Error: Template file '{args.template}' not found.")
        return
    if not os.path.exists(args.questions):
        print(f"❌ Error: Questions file '{args.questions}' not found.")
        return
    if args.image and not os.path.exists(args.image):
        print(f"❌ Error: Image file '{args.image}' not found.")
        return

    with open(args.questions, 'r', encoding='utf-8') as f:
        questions = json.load(f)

    prs = Presentation(args.template)
    
    print("--- Starting Presentation Generation ---")

    if args.image and len(prs.slides) > 0:
        print("\nProcessing Title Slide...")
        replace_title_image(prs.slides[0], args.image)

    print("\nProcessing Questions...")
    question_slides = []
    for slide in prs.slides:
        shape = find_question_shape(slide)
        if shape:
            question_slides.append((slide, shape))

    print(f"Found {len(question_slides)} question slides in template.")
    print(f"Loaded {len(questions)} questions from {args.questions}.")

    if len(questions) > len(question_slides):
        print(f"⚠  Warning: {len(questions)} questions but only {len(question_slides)} slides. Extra questions ignored.")

    for idx, question in enumerate(questions):
        if idx >= len(question_slides):
            break
        slide, shape = question_slides[idx]
        replace_question_content(shape, idx + 1, question)
        print(f"  ✓ Slide {idx+1} updated: {question['question_en'][:50]}...")

    prs.save(args.output)
    print(f"\n✅ Done! Presentation saved to: {args.output}")


if __name__ == '__main__':
    main()