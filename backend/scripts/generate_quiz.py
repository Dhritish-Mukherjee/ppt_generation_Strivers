"""
generate_quiz.py
----------------
Generates a quiz PPTX from slide_master.pptx and a questions.json file.

Slide structure produced:
  Slide 1    — Cover (always, copied from template slide 1)
  Slides 2-6 — Quiz_Layout × 5 questions
  Slide 7    — Promotional_Layout
  Slides 8-12— Quiz_Layout × next 5 questions
  Slide 13   — Promotional_Layout
  … and so on.
  The deck always ends with a Promotional_Layout slide.

Usage:
  python generate_quiz.py
  python generate_quiz.py --template slide_master.pptx --questions questions.json --output output.pptx

JSON format expected (same as before):
  [
    {
      "question_en": "Which of the following ...",
      "question_bn": "নিচের কোনটি ...",
      "options": [
        { "en": "A. Option text", "bn": "বাংলা টেক্সট" },
        { "en": "B. Option text", "bn": "বাংলা টেক্সট" },
        { "en": "C. Option text", "bn": "বাংলা টেক্সট" },
        { "en": "D. Option text", "bn": "বাংলা টেক্সট" }
      ]
    },
    ...
  ]
"""

import argparse
import copy
import json
import os
import re
import sys
from lxml import etree

# ── Namespace map used throughout ─────────────────────────────────────────
NS = {
    'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r':  'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'pkg':'http://schemas.openxmlformats.org/package/2006/relationships',
}

# ── Font size tiers (in hundredths of a point, i.e. "sz" attribute) ───────
# These mirror the original template values.
#   Question EN line: sz 4800 (~48pt)  → shrink to 3800 for long content
#   Question BN line: sz 3800 (~38pt)  → shrink to 3200 for long content
#   Option EN part:   sz 3909 (~39pt)  → shrink to 3200 for long content
#   Option BN part:   sz 2800 (~28pt)  → shrink to 2400 for long content

def _pick_sizes(question):
    """Return (q_en_sz, q_bn_sz, opt_en_sz, opt_bn_sz) based on content length."""
    all_text = (
        question['question_en'] + question['question_bn'] +
        ''.join(o['en'] + o.get('bn', '') for o in question['options'])
    )
    n = len(all_text)
    if n <= 350:
        return 4800, 3800, 3909, 2800
    elif n <= 550:
        return 4200, 3400, 3400, 2600
    else:
        return 3600, 3000, 3000, 2400


# ── XML helpers ────────────────────────────────────────────────────────────

def _a(tag):
    """Shorthand: return Clark-notation tag for the 'a' namespace."""
    return '{%s}%s' % (NS['a'], tag)


def _set_sz(rPr, sz):
    """Set the sz attribute on an <a:rPr> element."""
    rPr.set('sz', str(sz))


def _text_of(elem):
    """Concatenate all <a:t> descendants."""
    return ''.join(t.text or '' for t in elem.iter(_a('t')))


def _clear_runs(para):
    """Remove all <a:r> children from a paragraph, keeping <a:pPr> and <a:endParaRPr>."""
    for r in para.findall(_a('r')):
        para.remove(r)


def _first_rPr(para):
    """Return the first <a:rPr> found in a paragraph's runs (or endParaRPr), for cloning."""
    for r in para.findall(_a('r')):
        rPr = r.find(_a('rPr'))
        if rPr is not None:
            return rPr
    return para.find(_a('endParaRPr'))


def _make_run(rPr_template, lang, text, sz, bold=None):
    """
    Build a new <a:r> element cloned from rPr_template, with updated lang/sz/bold/text.
    """
    r = etree.Element(_a('r'))
    rPr = copy.deepcopy(rPr_template)
    rPr.tag = _a('rPr')
    rPr.set('lang', lang)
    _set_sz(rPr, sz)
    if bold is not None:
        rPr.set('b', '1' if bold else '0')
    # Remove err attribute if present (spellcheck artifact)
    if 'err' in rPr.attrib:
        del rPr.attrib['err']
    r.append(rPr)
    t = etree.SubElement(r, _a('t'))
    t.text = text
    if text and (text[0] == ' ' or text[-1] == ' '):
        t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    return r


# ── Core slide-filling logic ───────────────────────────────────────────────

def fill_quiz_slide(slide_xml_root, q_num, question):
    """
    Mutate slide_xml_root in-place to display the given question.

    Targets:
      - Shape named "Question_Fixed"  → 2 paragraphs: EN question, BN question
      - Shape named "Option_Fixed"    → 4 paragraphs: one per option (A-D)
    """
    q_en_sz, q_bn_sz, opt_en_sz, opt_bn_sz = _pick_sizes(question)

    spTree = slide_xml_root.find('.//' + '{%s}spTree' % NS['p'])
    if spTree is None:
        spTree = slide_xml_root.find('.//{%s}spTree' % NS['a'])

    for sp in slide_xml_root.iter('{%s}sp' % NS['p']):
        cNvPr = sp.find('.//{%s}cNvPr' % NS['p'])
        if cNvPr is None:
            continue
        name = cNvPr.get('name', '')

        if name == 'Question_Fixed':
            _fill_question(sp, q_num, question, q_en_sz, q_bn_sz)

        elif name == 'Option_Fixed':
            _fill_options(sp, question, opt_en_sz, opt_bn_sz)


def _fill_question(sp, q_num, question, q_en_sz, q_bn_sz):
    txBody = sp.find('{%s}txBody' % NS['p'])
    paras  = txBody.findall(_a('p'))

    # ── Paragraph 0: English question ─────────────────────────────────────
    p0 = paras[0]
    rPr0 = _first_rPr(p0) or etree.Element(_a('rPr'))
    _clear_runs(p0)
    en_text = f'{q_num}. {question["question_en"]}'
    p0.append(_make_run(rPr0, 'en-US', en_text, q_en_sz))

    # ── Paragraph 1: Bengali question ─────────────────────────────────────
    p1 = paras[1]
    rPr1 = _first_rPr(p1) or copy.deepcopy(rPr0)
    _clear_runs(p1)
    bn_text = question['question_bn']
    p1.append(_make_run(rPr1, 'hi-IN', bn_text, q_bn_sz))


def _fill_options(sp, question, opt_en_sz, opt_bn_sz):
    txBody = sp.find('{%s}txBody' % NS['p'])
    paras  = txBody.findall(_a('p'))
    letters = ['A', 'B', 'C', 'D']

    # Find the 4 option paragraphs by scanning for A./B./C./D. prefixes,
    # falling back to positional assignment.
    opt_paras = {}
    for p in paras:
        txt = _text_of(p).strip()
        for letter in letters:
            if txt.startswith(f'{letter}.'):
                opt_paras[letter] = p
                break

    # If detection failed (e.g. first run), assign positionally
    if len(opt_paras) < 4:
        content_paras = [p for p in paras if _text_of(p).strip()]
        for i, letter in enumerate(letters):
            if letter not in opt_paras and i < len(content_paras):
                opt_paras[letter] = content_paras[i]

    for i, opt in enumerate(question['options'][:4]):
        letter = letters[i]
        if letter not in opt_paras:
            continue

        p = opt_paras[letter]
        # Grab the rPr template from the first run (EN bold part)
        rPr_template = _first_rPr(p) or etree.Element(_a('rPr'))
        _clear_runs(p)

        # Clean option text (strip leading "A. " etc. if present)
        clean_en = re.sub(r'^[A-D]\.\s*', '', opt['en']).strip()
        clean_bn = re.sub(r'^[A-D]\.\s*', '', opt.get('bn', '')).strip()

        # Run 1: "A. English text " — bold, larger
        p.append(_make_run(rPr_template, 'en-US', f'{letter}. {clean_en} ', opt_en_sz, bold=True))

        # Run 2: "(" — smaller, not bold
        p.append(_make_run(rPr_template, 'en-US', '(', opt_bn_sz, bold=False))

        # Run 3: Bengali text — bold, smaller, Bengali locale
        if clean_bn:
            p.append(_make_run(rPr_template, 'hi-IN', clean_bn, opt_bn_sz, bold=True))

        # Run 4: ")" — bold, smaller
        p.append(_make_run(rPr_template, 'en-US', ')', opt_bn_sz, bold=True))


# ── PPTX manipulation via python-pptx ─────────────────────────────────────

def _build_slide_sequence(num_questions, batch_size=5):
    """
    Return a list of ('quiz', q_index) or ('promo',) tuples representing
    the desired slide order (excluding the cover slide).

    Pattern: [5× quiz, promo, 5× quiz, promo, …, always ends with promo]
    """
    seq = []
    for i in range(num_questions):
        seq.append(('quiz', i))
        if (i + 1) % batch_size == 0:
            seq.append(('promo',))
    # Always end with a promo slide (unless the last slide was already a promo)
    if not seq or seq[-1][0] != 'promo':
        seq.append(('promo',))
    return seq


def main():
    parser = argparse.ArgumentParser(description='Generate quiz PPTX from slide master')
    parser.add_argument('--template',  default='slide_master.pptx', help='Path to slide_master.pptx')
    parser.add_argument('--questions', default='questions.json',    help='Path to questions JSON')
    parser.add_argument('--output',    default='output.pptx',       help='Output file name')
    parser.add_argument('--batch',     default=5, type=int,         help='Questions per batch before a promo slide (default 5)')
    args = parser.parse_args()

    # ── Validate inputs ────────────────────────────────────────────────────
    for path, label in [(args.template, 'Template'), (args.questions, 'Questions')]:
        if not os.path.exists(path):
            print(f'❌ {label} file not found: {path}')
            sys.exit(1)

    with open(args.questions, 'r', encoding='utf-8') as f:
        questions = json.load(f)

    print(f'Loaded {len(questions)} question(s) from {args.questions}')
    print(f'Batch size: {args.batch} questions per promo slide')

    # ── Use python-pptx to clone slides ───────────────────────────────────
    from pptx import Presentation
    from pptx.util import Emu
    import copy as _copy
    from lxml import etree as _etree

    prs = Presentation(args.template)
    slides = prs.slides

    # Identify source slides
    # slide index 0 = cover, 1 = quiz template, 2 = promo template
    if len(slides) < 3:
        print('❌ Template must have at least 3 slides (cover, quiz, promo).')
        sys.exit(1)

    cover_slide = slides[0]
    quiz_slide  = slides[1]
    promo_slide = slides[2]

    # ── Build a fresh presentation copying the slide master/layout refs ───
    # Strategy: start from the template, delete slides 1-2 (quiz+promo),
    # then re-add them in the desired pattern.
    # python-pptx doesn't support slide deletion natively, so we work at
    # the XML level via the presentation.xml sldIdLst.

    prs_xml = prs.element

    # Grab references to the two source slides' XML so we can deep-copy them
    def _slide_xml(slide):
        return slide._element

    quiz_xml_orig  = _slide_xml(quiz_slide)
    promo_xml_orig = _slide_xml(promo_slide)

    # We'll use python-pptx's add_slide with the correct layout
    quiz_layout  = quiz_slide.slide_layout
    promo_layout = promo_slide.slide_layout

    # Remove the template quiz+promo slides from the deck (keep cover only)
    # We do this by manipulating the slide list at XML level
    slide_id_lst = prs_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')

    # Get rIds of slides we want to remove (quiz and promo placeholders)
    def _get_slide_rId(slide):
        for rel in prs.part.rels.values():
            if rel._target is slide.part:
                return rel.rId
        return None

    quiz_rId  = _get_slide_rId(quiz_slide)
    promo_rId = _get_slide_rId(promo_slide)

    # Build the sequence of slides to add
    sequence = _build_slide_sequence(len(questions), args.batch)
    print(f'\nSlide plan: 1 cover + {len(sequence)} content slides')
    print(f'  ({sum(1 for s in sequence if s[0]=="quiz")} quiz slides, '
          f'{sum(1 for s in sequence if s[0]=="promo")} promo slides)\n')

    # Add all needed slides using add_slide, then fill quiz ones with content
    new_slides = []
    for item in sequence:
        if item[0] == 'quiz':
            new_slide = prs.slides.add_slide(quiz_layout)
            # Deep-copy the body content from the original quiz slide XML
            # (add_slide creates a blank layout-based slide; we copy the sp tree)
            orig_spTree = quiz_xml_orig.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
            new_spTree  = new_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
            if orig_spTree is not None and new_spTree is not None:
                # Replace new spTree with deep copy of original
                parent = new_spTree.getparent()
                idx    = list(parent).index(new_spTree)
                parent.remove(new_spTree)
                parent.insert(idx, copy.deepcopy(orig_spTree))
            new_slides.append(('quiz', item[1], new_slide))
        else:
            new_slide = prs.slides.add_slide(promo_layout)
            orig_spTree = promo_xml_orig.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
            new_spTree  = new_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
            if orig_spTree is not None and new_spTree is not None:
                parent = new_spTree.getparent()
                idx    = list(parent).index(new_spTree)
                parent.remove(new_spTree)
                parent.insert(idx, copy.deepcopy(orig_spTree))
            new_slides.append(('promo', None, new_slide))

    # Now remove the original template quiz+promo slides from sldIdLst
    # (they are slides at index 1 and 2 in the original deck)
    sldIdLst = prs.element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst'
    )
    all_sld_ids = list(sldIdLst)
    # Index 0 = cover, 1 = quiz template, 2 = promo template → remove 1 and 2
    for sld_id_elem in all_sld_ids[1:3]:
        sldIdLst.remove(sld_id_elem)

    print('Filling slides with question content...')

    # Fill quiz slides with question data
    for kind, q_idx, slide in new_slides:
        if kind == 'quiz':
            q = questions[q_idx]
            fill_quiz_slide(slide._element, q_idx + 1, q)
            print(f'  ✓ Q{q_idx+1}: {q["question_en"][:55]}...')
        else:
            print(f'  ── Promotional slide')

    prs.save(args.output)
    print(f'\n✅ Saved: {args.output}')
    total = 1 + len(new_slides)  # cover + content
    print(f'   Total slides: {total}')


if __name__ == '__main__':
    main()